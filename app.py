import os
import io
import json
import logging
from flask import Flask, request, render_template, jsonify, session
from werkzeug.utils import secure_filename
import google.generativeai as genai
from PIL import Image
import PyPDF2
from pptx import Presentation
import uuid
import traceback
from flask_cors import CORS
import pickle
import tempfile
from datetime import datetime, timedelta

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

app = Flask(__name__)
# Use a fixed secret key for session persistence
app.secret_key = 'your-fixed-secret-key-here-change-this-in-production'
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SECURE'] = False  # Set to True in production with HTTPS
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['PERMANENT_SESSION_LIFETIME'] = 7200  # 2 hours

# Enable CORS for frontend-backend communication
CORS(app, supports_credentials=True)

# Configure Gemini AI
genai.configure(api_key=os.getenv('GEMINI_API_KEY'))
model = genai.GenerativeModel('gemini-2.0-flash-exp')

# Upload configuration
UPLOAD_FOLDER = 'uploads'
SESSION_DATA_FOLDER = 'session_data'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 'ppt', 'pptx', 'doc', 'docx'}
MAX_CONTENT_LENGTH = 100 * 1024 * 1024  # 100MB

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# Ensure directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SESSION_DATA_FOLDER, exist_ok=True)

# In-memory storage for session data as backup
session_storage = {}

def save_session_data(session_id, data):
    """Save session data to both file and memory"""
    try:
        # Save to file
        file_path = os.path.join(SESSION_DATA_FOLDER, f"{session_id}.pkl")
        with open(file_path, 'wb') as f:
            pickle.dump(data, f)
        
        # Save to memory as backup
        session_storage[session_id] = {
            'data': data,
            'timestamp': datetime.now()
        }
        
        logger.info(f"Session data saved for ID: {session_id}")
        return True
    except Exception as e:
        logger.error(f"Error saving session data: {e}")
        # Still save to memory
        session_storage[session_id] = {
            'data': data,
            'timestamp': datetime.now()
        }
        return False

def load_session_data(session_id):
    """Load session data from file or memory"""
    try:
        # Try loading from file first
        file_path = os.path.join(SESSION_DATA_FOLDER, f"{session_id}.pkl")
        if os.path.exists(file_path):
            with open(file_path, 'rb') as f:
                data = pickle.load(f)
                logger.info(f"Session data loaded from file for ID: {session_id}")
                return data
    except Exception as e:
        logger.error(f"Error loading session data from file: {e}")
    
    # Try loading from memory
    if session_id in session_storage:
        stored_data = session_storage[session_id]
        # Check if data is not too old (2 hours)
        if datetime.now() - stored_data['timestamp'] < timedelta(hours=2):
            logger.info(f"Session data loaded from memory for ID: {session_id}")
            return stored_data['data']
        else:
            # Remove old data
            del session_storage[session_id]
    
    logger.warning(f"No session data found for ID: {session_id}")
    return None

def cleanup_old_sessions():
    """Clean up old session files and memory data"""
    try:
        # Clean up memory
        cutoff_time = datetime.now() - timedelta(hours=2)
        old_keys = [k for k, v in session_storage.items() if v['timestamp'] < cutoff_time]
        for key in old_keys:
            del session_storage[key]
        
        # Clean up files
        for filename in os.listdir(SESSION_DATA_FOLDER):
            if filename.endswith('.pkl'):
                file_path = os.path.join(SESSION_DATA_FOLDER, filename)
                file_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                if datetime.now() - file_time > timedelta(hours=2):
                    os.remove(file_path)
    except Exception as e:
        logger.error(f"Error cleaning up old sessions: {e}")

def get_or_create_session_id():
    """Get existing session ID or create a new one"""
    if 'session_id' not in session:
        session['session_id'] = str(uuid.uuid4())
        session.permanent = True
    return session['session_id']

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                text += f"\n--- Page {page_num} ---\n{page_text}\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""

def extract_text_from_ppt(file_path):
    """Extract text from PowerPoint file"""
    try:
        text = ""
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += f"{shape.text}\n"
                # Extract text from tables if present
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text += f"Table: {' | '.join(row_text)}\n"
            text += "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PPT text: {e}")
        return ""

def process_image(file_path):
    """Process image using Gemini Vision"""
    try:
        img = Image.open(file_path)
        prompt = """
        Analyze this image thoroughly and extract ALL educational content. Be extremely detailed and include:
        
        1. **All visible text** - transcribe every word, sentence, heading, caption, and label visible in the image
        2. **Mathematical formulas and equations** - write out every formula, equation, or mathematical expression exactly as shown
        3. **Diagrams and visual elements** - describe every diagram, chart, graph, flowchart, or visual representation in detail
        4. **Tables and data** - extract all tabular data, measurements, values, and statistics
        5. **Scientific concepts** - identify and explain any scientific principles, theories, or concepts shown
        6. **Technical details** - include specifications, dimensions, measurements, or technical parameters
        7. **Relationships and connections** - describe how different elements in the image relate to each other
        8. **Context and explanations** - provide context for what each element represents or explains
        
        Format your response clearly with proper headings and organize the information logically. Don't miss any detail, no matter how small - this content will be used to create comprehensive learning materials.
        """
        
        response = model.generate_content([prompt, img])
        return response.text
    except Exception as e:
        logger.error(f"Error processing image: {e}")
        return ""

def generate_learning_content(content, file_type="unknown"):
    """Generate comprehensive learning content using Gemini"""
    prompt = f"""
    Based on the following educational content from a {file_type} file, create a comprehensive learning package. The content is:

    Content: {content[:12000]}  # Increased limit for more detailed processing

    Please provide a JSON response with the following structure:
    {{
        "summary": "Brief but comprehensive overview of the main topic and all key concepts covered",
        "detailed_explanation": "In-depth explanation of all topics, formulas, concepts, and principles covered in the content. Include step-by-step explanations, examples, and context.",
        "flashcards": [
            {{
                "front": "Clear, specific question or concept to test understanding",
                "back": "Detailed, comprehensive answer that includes: full explanations, all relevant formulas with proper notation, step-by-step processes, examples where applicable, key points and important details. Make each answer substantial and educational - aim for 150-300 words per flashcard back to ensure complete understanding."
            }},
            // Create 20-25 flashcards covering ALL concepts, formulas, definitions, processes, and important details from the content
        ],
        "quiz_questions": [
            {{
                "question": "Well-crafted question testing understanding of concepts",
                "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
                "correct_answer": "A",
                "explanation": "Detailed explanation of why this answer is correct and why others are wrong, including relevant formulas or concepts"
            }},
            // Include exactly 15 questions covering all major topics and concepts
        ],
        "key_topics": ["List all important topics, concepts, formulas, and principles covered"],
        "learning_objectives": ["Specific, measurable learning goals covering all content areas"],
        "important_formulas": ["List all mathematical formulas, equations, and expressions found in the content with proper notation"],
        "diagrams_explained": ["Detailed explanations of any diagrams, charts, or visual elements mentioned in the content"]
    }}

    CRITICAL REQUIREMENTS:
    - Make flashcard answers VERY detailed and comprehensive (150-300 words each)
    - Include ALL formulas, equations, and mathematical expressions exactly as they appear
    - Cover EVERY concept, definition, and important point from the content
    - Provide step-by-step explanations for processes and procedures
    - Include examples and applications where relevant
    - Don't skip any topic or detail, no matter how small
    - Make explanations clear and educational but thorough
    - Ensure flashcards can serve as complete study guides for each topic
    """

    try:
        response = model.generate_content(prompt)
        
        # Try to extract JSON from the response
        response_text = response.text.strip()
        
        # Remove markdown code blocks if present
        if response_text.startswith('```json'):
            response_text = response_text[7:]
        if response_text.endswith('```'):
            response_text = response_text[:-3]
        
        # Parse JSON
        learning_data = json.loads(response_text)
        
        # Validate structure
        required_keys = ['summary', 'detailed_explanation', 'flashcards', 'quiz_questions', 'key_topics']
        for key in required_keys:
            if key not in learning_data:
                raise ValueError(f"Missing required key: {key}")
        
        # Ensure we have enough flashcards and they are detailed
        if len(learning_data['flashcards']) < 15:
            logger.warning("Generated fewer flashcards than expected")
        
        return learning_data
        
    except json.JSONDecodeError as e:
        logger.error(f"JSON parsing error: {e}")
        return create_enhanced_fallback_content(content)
    except Exception as e:
        logger.error(f"Error generating learning content: {e}")
        return create_enhanced_fallback_content(content)

def create_enhanced_fallback_content(content):
    """Create enhanced fallback content if AI generation fails"""
    # Extract some basic information from content
    content_preview = content[:2000] if len(content) > 2000 else content
    
    return {
        "summary": f"Educational content has been uploaded and processed. The material covers various topics and concepts that require detailed study.",
        "detailed_explanation": content_preview + "...\n\nThis appears to be educational material that covers multiple concepts. Please try regenerating the learning materials for a more detailed analysis.",
        "flashcards": [
            {
                "front": "What is the main subject matter of this educational content?",
                "back": f"This educational material covers various concepts and topics. Based on the content preview: {content_preview[:500]}... The material appears to contain detailed information that requires careful study and understanding. To get the most out of this content, review it section by section and identify key concepts, formulas, and important principles."
            },
            {
                "front": "What should be the approach to studying this material?",
                "back": "To effectively study this material: 1) Read through the entire content to get an overview, 2) Identify key concepts and important points, 3) Pay special attention to any formulas, definitions, or technical details, 4) Practice applying the concepts through examples, 5) Review regularly to reinforce understanding. Take notes on important points and create your own summaries of complex topics."
            }
        ],
        "quiz_questions": [
            {
                "question": "What is the recommended approach when studying new educational material?",
                "options": [
                    "A) Read once and move on",
                    "B) Focus only on memorizing facts",
                    "C) Review systematically and identify key concepts",
                    "D) Skip difficult sections"
                ],
                "correct_answer": "C",
                "explanation": "Systematic review and identification of key concepts ensures comprehensive understanding and better retention of the material."
            }
        ],
        "key_topics": ["Educational Content Analysis", "Study Methods", "Content Review"],
        "learning_objectives": ["Understand the uploaded content", "Identify key concepts", "Develop effective study strategies"],
        "important_formulas": ["Content contains formulas - please regenerate for detailed extraction"],
        "diagrams_explained": ["Visual elements present - please regenerate for detailed descriptions"]
    }

def chat_with_content(question, content_context, learning_data):
    """Enhanced chatbot that can answer questions about the content"""
    # Prepare comprehensive context including learning data
    context_info = f"""
    EDUCATIONAL CONTENT CONTEXT:
    {content_context[:4000]}
    
    LEARNING SUMMARY:
    {learning_data.get('summary', 'No summary available')}
    
    KEY TOPICS:
    {', '.join(learning_data.get('key_topics', []))}
    
    IMPORTANT FORMULAS:
    {', '.join(learning_data.get('important_formulas', []))}
    """
    
    prompt = f"""
    You are an expert AI tutor with deep knowledge of the educational content that the student has uploaded. 
    
    CONTEXT: {context_info}
    
    STUDENT QUESTION: {question}
    
    Please provide a comprehensive, helpful answer that:
    1. Directly addresses the student's question
    2. References specific information from their uploaded content
    3. Provides detailed explanations with examples when relevant
    4. Includes any relevant formulas or concepts from the material
    5. Uses clear, educational language that promotes understanding
    6. Encourages further learning and exploration of the topic
    
    If the question is about a specific formula, concept, or topic from the content, provide step-by-step explanations and practical applications. If the question is more general, connect it to the learning material whenever possible.
    
    Be encouraging, thorough, and educational in your response.
    """
    
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        logger.error(f"Chat error: {e}")
        return "I'm sorry, I'm having trouble processing your question right now. Could you please try rephrasing it or asking about a specific topic from your uploaded content?"

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Please upload a file smaller than 100MB.'}), 413

@app.errorhandler(Exception)
def handle_exception(e):
    logger.error(f"Unhandled exception: {e}")
    logger.error(traceback.format_exc())
    return jsonify({'error': 'An unexpected error occurred. Please try again.'}), 500

@app.route('/')
def index():
    # Clean up old sessions periodically
    cleanup_old_sessions()
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        # Get or create session ID
        session_id = get_or_create_session_id()
        
        # Check if file is present in request
        if 'file' not in request.files:
            logger.error("No file in request")
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            logger.error("No file selected")
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename):
            logger.error(f"File type not supported: {file.filename}")
            return jsonify({'error': 'File type not supported. Please upload PDF, PPT, PPTX, images, or text files.'}), 400
        
        logger.info(f"Processing file: {file.filename} for session: {session_id}")
        
        # Save file
        filename = secure_filename(file.filename)
        unique_filename = f"{uuid.uuid4()}_{filename}"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        
        try:
            file.save(file_path)
            logger.info(f"File saved to: {file_path}")
        except Exception as save_error:
            logger.error(f"Error saving file: {save_error}")
            return jsonify({'error': 'Failed to save file'}), 500
        
        # Extract content based on file type
        content = ""
        file_ext = filename.rsplit('.', 1)[1].lower()
        
        logger.info(f"Extracting content from {file_ext} file")
        
        try:
            if file_ext == 'pdf':
                content = extract_text_from_pdf(file_path)
            elif file_ext in ['ppt', 'pptx']:
                content = extract_text_from_ppt(file_path)
            elif file_ext in ['png', 'jpg', 'jpeg', 'gif']:
                content = process_image(file_path)
            elif file_ext == 'txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            else:
                logger.error(f"Unsupported file extension: {file_ext}")
                return jsonify({'error': 'Unsupported file type'}), 400
                
        except Exception as extract_error:
            logger.error(f"Error extracting content: {extract_error}")
            content = ""
        finally:
            # Clean up uploaded file
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    logger.info("Uploaded file cleaned up")
            except Exception as cleanup_error:
                logger.error(f"Failed to cleanup file: {cleanup_error}")
        
        if not content.strip():
            logger.error("No content extracted from file")
            return jsonify({'error': 'Could not extract readable content from the file. Please ensure the file contains text or readable content.'}), 400
        
        logger.info(f"Content extracted successfully. Length: {len(content)} characters")
        
        # Generate learning content
        try:
            learning_data = generate_learning_content(content, file_ext)
            logger.info("Learning content generated successfully")
        except Exception as gen_error:
            logger.error(f"Error generating learning content: {gen_error}")
            learning_data = create_enhanced_fallback_content(content)
        
        # Prepare session data
        session_data = {
            'content_context': content,
            'learning_data': learning_data,
            'original_filename': filename,
            'file_uploaded': True,
            'upload_timestamp': datetime.now().isoformat(),
            'file_type': file_ext
        }
        
        # Save session data using our robust system
        save_session_data(session_id, session_data)
        
        logger.info(f"Session data saved for session: {session_id}")
        
        return jsonify({
            'success': True,
            'data': learning_data,
            'filename': filename,
            'session_id': session_id
        })
        
    except Exception as e:
        logger.error(f"Upload error: {e}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'Processing failed: {str(e)}'}), 500

@app.route('/chat', methods=['POST'])
def chat():
    try:
        # Get session ID
        session_id = get_or_create_session_id()
        
        data = request.get_json()
        question = data.get('question', '')
        
        if not question.strip():
            return jsonify({'error': 'No question provided'}), 400
        
        logger.info(f"Chat request received for session: {session_id} - Question: {question[:100]}...")
        
        # Load session data using our robust system
        session_data = load_session_data(session_id)
        
        if not session_data:
            logger.error(f"No session data found for session: {session_id}")
            return jsonify({'error': 'No content context available. Please upload a file first.'}), 400
        
        # Check if we have the required data
        if not session_data.get('file_uploaded', False):
            logger.error(f"File not uploaded flag missing for session: {session_id}")
            return jsonify({'error': 'Please upload a file first before asking questions.'}), 400
        
        content_context = session_data.get('content_context', '')
        learning_data = session_data.get('learning_data', {})
        
        if not content_context:
            logger.error(f"Content context is empty for session: {session_id}")
            return jsonify({'error': 'Content not found. Please try uploading the file again.'}), 400
        
        logger.info(f"Using content context of length: {len(content_context)} for session: {session_id}")
        
        response = chat_with_content(question, content_context, learning_data)
        
        logger.info(f"Chat response generated successfully for session: {session_id}")
        
        return jsonify({
            'success': True,
            'response': response,
            'debug': {
                'content_length': len(content_context),
                'has_learning_data': bool(learning_data),
                'session_id': session_id,
                'upload_time': session_data.get('upload_timestamp', 'unknown')
            }
        })
        
    except Exception as e:
        logger.error(f"Chat error: {e}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'Chat processing failed: {str(e)}'}), 500

@app.route('/regenerate', methods=['POST'])
def regenerate_content():
    try:
        session_id = get_or_create_session_id()
        session_data = load_session_data(session_id)
        
        if not session_data:
            return jsonify({'error': 'No content available to regenerate. Please upload a file first.'}), 400
        
        content_context = session_data.get('content_context', '')
        if not content_context:
            return jsonify({'error': 'No content available to regenerate'}), 400
        
        # Get original file type
        file_ext = session_data.get('file_type', 'unknown')
        
        # Regenerate learning content
        learning_data = generate_learning_content(content_context, file_ext)
        
        # Update session data
        session_data['learning_data'] = learning_data
        save_session_data(session_id, session_data)
        
        return jsonify({
            'success': True,
            'data': learning_data
        })
        
    except Exception as e:
        logger.error(f"Regenerate error: {e}")
        return jsonify({'error': 'Regeneration failed. Please try again.'}), 500

# Add a debug route to check session status
@app.route('/debug/session', methods=['GET'])
def debug_session():
    session_id = get_or_create_session_id()
    session_data = load_session_data(session_id)
    
    if session_data:
        return jsonify({
            'session_id': session_id,
            'has_session_data': True,
            'file_uploaded': session_data.get('file_uploaded', False),
            'has_content_context': bool(session_data.get('content_context', '')),
            'has_learning_data': bool(session_data.get('learning_data', {})),
            'content_length': len(session_data.get('content_context', '')),
            'upload_timestamp': session_data.get('upload_timestamp', 'None'),
            'filename': session_data.get('original_filename', 'None'),
            'file_type': session_data.get('file_type', 'unknown')
        })
    else:
        return jsonify({
            'session_id': session_id,
            'has_session_data': False,
            'message': 'No session data found'
        })

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)